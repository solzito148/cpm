#!/usr/bin/env python3
"""Generate Cashflow weekly status PPTX from CPM-ProcurementForecasting-Status.xlsx."""

from __future__ import annotations

import re
from datetime import date, datetime
from pathlib import Path

import openpyxl
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

NAVY = RGBColor(0x1B, 0x2A, 0x4A)
BLUE = RGBColor(0x34, 0x83, 0xFA)
CYAN = RGBColor(0x00, 0xBC, 0xD4)
GREEN = RGBColor(0x00, 0xA6, 0x50)
PURPLE = RGBColor(0x9C, 0x7C, 0xDB)
RED = RGBColor(0xF2, 0x3D, 0x4F)
ORANGE = RGBColor(0xFF, 0x77, 0x33)
GRAY = RGBColor(0x9E, 0x9E, 0x9E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_LIGHT = RGBColor(0xEE, 0xF2, 0xF7)
TEXT_PRI = RGBColor(0x1B, 0x2A, 0x4A)
TEXT_SEC = RGBColor(0x66, 0x66, 0x66)
SLATE = RGBColor(0x78, 0x90, 0x9C)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

STATUS_ORDER = [
    "Done",
    "Ready to move to RPC",
    "In Development",
    "In Progress",
    "In analysis",
    "To Do",
]

STATUS_COLORS = {
    "Done": GREEN,
    "Ready to move to RPC": PURPLE,
    "Ready for deployment": PURPLE,
    "In Development": BLUE,
    "In Progress": BLUE,
    "In analysis": CYAN,
    "To Do": GRAY,
    "Testing": PURPLE,
}

IN_PROGRESS_STATUSES = {"In Development", "In Progress", "In analysis"}

CONTENT_LEFT = Inches(0.4)
CONTENT_WIDTH = Inches(12.55)
TABLE_TOP = Inches(1.05)
FONT = "Segoe UI"


def truncate_words(text: str, max_len: int) -> str:
    text = clean(text)
    if len(text) <= max_len:
        return text
    cut = text[: max_len - 1].rsplit(" ", 1)[0]
    return (cut or text[: max_len - 1]).rstrip(",.;") + "…"


def style_cell(
    cell,
    *,
    font_size: int = 9,
    bold: bool = False,
    color=TEXT_PRI,
    align=PP_ALIGN.LEFT,
    wrap: bool = True,
):
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Pt(4)
    cell.margin_right = Pt(4)
    cell.margin_top = Pt(2)
    cell.margin_bottom = Pt(2)
    tf = cell.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = color
        paragraph.font.name = FONT
        paragraph.alignment = align
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(0)


def set_cell_text(cell, text: str, **kwargs):
    cell.text = str(text)
    style_cell(cell, **kwargs)


def add_multiline_text_box(
    slide,
    left,
    top,
    width,
    height,
    lines: list[str],
    font_size=12,
    bold=False,
    color=TEXT_PRI,
    align=PP_ALIGN.LEFT,
    line_spacing=1.15,
):
    tx_box = slide.shapes.add_textbox(left, top, width, height)
    tf = tx_box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = FONT
        p.alignment = align
        p.line_spacing = line_spacing
    return tx_box


def clean(value) -> str:
    if value is None:
        return ""
    text = str(value).replace("\xa0", " ").strip()
    return re.sub(r"\s+", " ", text)


def fmt_date(value) -> str:
    if value in (None, "", "TBD"):
        return "TBD"
    if isinstance(value, datetime):
        return value.strftime("%d-%b")
    if isinstance(value, date):
        return value.strftime("%d-%b")
    return clean(value)


def fmt_priority(value) -> str:
    if value in (None, "", "-"):
        return "-"
    return str(value)


def sci_code(raw: str) -> str:
    text = clean(raw).upper().replace("SCI-", "SCI-")
    match = re.search(r"SCI[-\s]?(\d+)", text, re.I)
    if not match:
        return clean(raw)
    return f"SCI-{int(match.group(1)):02d}"


def sci_sort_key(code: str) -> tuple[int, str]:
    match = re.search(r"SCI-(\d+)", code)
    return (int(match.group(1)) if match else 999, code)


def load_tasks(xlsx_path: Path) -> list[dict]:
    wb = openpyxl.load_workbook(xlsx_path, data_only=True)
    ws = wb["Backlog"]
    tasks: list[dict] = []

    for row in ws.iter_rows(min_row=2, values_only=True):
        task_type, ritm, item_id, jira, name, status = row[:6]
        assigned, requestor, priority, due, initial, deployment, change_req = row[6:13]
        details = row[13] if len(row) > 13 else ""

        if not any(clean(v) for v in (jira, name, status)):
            continue

        code = sci_code(jira)
        if not code.startswith("SCI-"):
            continue

        status_text = clean(status)
        tasks.append(
            {
                "type": clean(task_type) or "Enhancement",
                "ritm": clean(ritm) or "-",
                "id": clean(item_id) or "-",
                "code": code,
                "name": clean(name),
                "status": status_text,
                "assigned": clean(assigned) or "TBD",
                "requestor": clean(requestor) or "-",
                "priority": fmt_priority(priority),
                "due": fmt_date(due),
                "initial": fmt_date(initial),
                "deployment": fmt_date(deployment),
                "change_request": clean(change_req) or "-",
                "details": clean(details),
            }
        )

    tasks.sort(key=lambda t: sci_sort_key(t["code"]))
    return tasks


def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_oval(slide, left, top, width, height, fill_color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if alpha is not None:
        from lxml import etree
        from pptx.oxml.ns import qn

        sp_pr = shape._element.spPr
        solid_fill = sp_pr.find(qn("a:solidFill"))
        if solid_fill is not None:
            srgb_clr = solid_fill.find(qn("a:srgbClr"))
            if srgb_clr is not None:
                srgb_clr.append(srgb_clr.makeelement(qn("a:alpha"), {"val": str(alpha)}))
    return shape


def add_text_box(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size=14,
    bold=False,
    color=TEXT_PRI,
    align=PP_ALIGN.LEFT,
):
    tx_box = slide.shapes.add_textbox(left, top, width, height)
    tx_box.text_frame.word_wrap = True
    p = tx_box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = FONT
    p.alignment = align
    tf = tx_box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    return tx_box


def add_header_bar(slide):
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.9), NAVY)
    add_text_box(slide, Inches(11), Inches(0.22), Inches(2.1), Inches(0.5),
                 "Cashflow Project", 16, False, CYAN, PP_ALIGN.RIGHT)


def set_header_title(slide, title):
    font_size = 20 if len(title) > 34 else 24
    add_text_box(
        slide,
        Inches(0.45),
        Inches(0.16),
        Inches(10.2),
        Inches(0.65),
        title,
        font_size,
        True,
        WHITE,
    )


LOGO_PATH = Path(__file__).resolve().parent / "assets" / "aes-logo.png"


def add_aes_logo(slide, left=Inches(8.15), top=Inches(6.05), width=Inches(3.0)):
    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), left, top, width=width)
        return

    tx_logo = slide.shapes.add_textbox(left, top, Inches(2.5), Inches(1))
    p = tx_logo.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    for char, clr in [("a", BLUE), ("e", PURPLE), ("s", GREEN)]:
        run = p.add_run()
        run.text = char
        run.font.size = Pt(56)
        run.font.bold = True
        run.font.color.rgb = clr
        run.font.name = "Segoe UI"


def create_table_slide(slide, title, rows, cols_spec, *, row_height=Inches(0.38), font_size=9):
    add_bg(slide, BG_LIGHT)
    add_header_bar(slide)
    set_header_title(slide, title)

    n_cols = len(cols_spec)
    n_rows = len(rows) + 1
    tbl_left = CONTENT_LEFT
    tbl_top = TABLE_TOP
    tbl_w = sum(width for _, width in cols_spec)
    tbl_h = row_height * n_rows

    table_shape = slide.shapes.add_table(n_rows, n_cols, tbl_left, tbl_top, tbl_w, tbl_h)
    table = table_shape.table

    for ci, (_, width) in enumerate(cols_spec):
        table.columns[ci].width = width

    for ci, (header, _) in enumerate(cols_spec):
        cell = table.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        set_cell_text(cell, header, font_size=10, bold=True, color=WHITE)

    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
            header = cols_spec[ci][0]
            color = TEXT_PRI
            bold = False
            if header == "SCI":
                color = CYAN
                bold = True
            elif header == "Status":
                color = STATUS_COLORS.get(str(val), TEXT_PRI)
                bold = True
            elif header == "Due" and str(val) not in ("-", "TBD", ""):
                color = ORANGE
                bold = True
            set_cell_text(cell, val, font_size=font_size, bold=bold, color=color)


def detail_card(slide, x, y, task, accent=BLUE):
    card_w = Inches(6.05)
    card_h = Inches(1.62)
    add_rect(slide, x, y, card_w, card_h, WHITE)
    add_rect(slide, x, y, Inches(0.07), card_h, accent)

    add_rounded_rect(slide, x + Inches(0.14), y + Inches(0.12), Inches(0.78), Inches(0.3), CYAN)
    add_text_box(
        slide,
        x + Inches(0.14),
        y + Inches(0.12),
        Inches(0.78),
        Inches(0.3),
        task["code"],
        10,
        True,
        WHITE,
        PP_ALIGN.CENTER,
    )

    meta_x = x + Inches(0.98)
    if task["priority"] not in ("-", ""):
        badge_color = RED if task["priority"] == "1" else BLUE
        add_rounded_rect(slide, meta_x, y + Inches(0.12), Inches(0.38), Inches(0.3), badge_color)
        add_text_box(
            slide,
            meta_x,
            y + Inches(0.12),
            Inches(0.38),
            Inches(0.3),
            f"P{task['priority']}",
            9,
            True,
            WHITE,
            PP_ALIGN.CENTER,
        )
        meta_x += Inches(0.45)

    owner = task["requestor"] if task["requestor"] != "-" else "—"
    add_text_box(
        slide,
        x + Inches(2.55),
        y + Inches(0.12),
        Inches(3.35),
        Inches(0.32),
        f"Requestor: {owner} · Assigned: {task['assigned']}",
        8,
        False,
        TEXT_SEC,
        PP_ALIGN.RIGHT,
    )

    add_text_box(
        slide,
        x + Inches(0.14),
        y + Inches(0.48),
        Inches(5.75),
        Inches(0.55),
        truncate_words(task["name"], 95),
        11,
        True,
        TEXT_PRI,
    )

    detail = task["details"] or task["ritm"]
    add_text_box(
        slide,
        x + Inches(0.14),
        y + Inches(0.98),
        Inches(5.75),
        Inches(0.38),
        truncate_words(detail, 110),
        8,
        False,
        TEXT_SEC,
    )

    add_text_box(
        slide,
        x + Inches(0.14),
        y + Inches(1.28),
        Inches(2.4),
        Inches(0.24),
        task["status"],
        9,
        True,
        STATUS_COLORS.get(task["status"], TEXT_PRI),
    )

    due_text = task["due"] if task["due"] not in ("TBD", "-", "") else "No due date"
    due_color = ORANGE if due_text != "No due date" else GRAY
    add_text_box(
        slide,
        x + Inches(4.35),
        y + Inches(1.28),
        Inches(1.55),
        Inches(0.24),
        due_text,
        9,
        True,
        due_color,
        PP_ALIGN.RIGHT,
    )


def chunk(items, size):
    for i in range(0, len(items), size):
        yield items[i : i + size]


ROWS_PER_TABLE_PAGE = 10


def table_page_suffix(page_idx: int, total_pages: int) -> str:
    return f" ({page_idx}/{total_pages})" if total_pages > 1 else ""


def add_paginated_table_slides(
    prs,
    blank,
    title: str,
    rows: list,
    cols_spec,
    *,
    row_height=Inches(0.42),
    font_size=9,
):
    total_pages = max(1, (len(rows) + ROWS_PER_TABLE_PAGE - 1) // ROWS_PER_TABLE_PAGE)
    for page_idx, page_rows in enumerate(chunk(rows, ROWS_PER_TABLE_PAGE), start=1):
        slide = prs.slides.add_slide(blank)
        create_table_slide(
            slide,
            f"{title}{table_page_suffix(page_idx, total_pages)}",
            page_rows,
            cols_spec,
            row_height=row_height,
            font_size=font_size,
        )


def build_presentation(tasks: list[dict], report_date: date, output_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    ready = [t for t in tasks if t["status"] == "Ready to move to RPC"]
    active = [t for t in tasks if t["status"] in IN_PROGRESS_STATUSES]
    todo = [t for t in tasks if t["status"] == "To Do"]
    done = [t for t in tasks if t["status"] == "Done"]
    priority_items = [t for t in tasks if t["priority"] == "1"]

    # Slide 1 — title
    slide1 = prs.slides.add_slide(blank)
    add_bg(slide1, WHITE)
    add_rect(slide1, Inches(0), Inches(0), Inches(0.08), SLIDE_H, CYAN)
    add_oval(slide1, Inches(9.8), Inches(-1), Inches(4.8), Inches(4.8), RGBColor(0xD6, 0xE4, 0xF7), 30000)
    add_oval(slide1, Inches(10.5), Inches(5.2), Inches(3.5), Inches(3.5), RGBColor(0xD0, 0xF0, 0xF4), 20000)
    add_text_box(slide1, Inches(1), Inches(2), Inches(8), Inches(1.2), "CASHFLOW PROJECT", 52, True, NAVY)
    add_text_box(slide1, Inches(1), Inches(3.4), Inches(8), Inches(0.5),
                 f"Status Update  |  {report_date.strftime('%B %d, %Y')}", 22, False, BLUE)
    add_rect(slide1, Inches(1), Inches(4.05), Inches(5.5), Inches(0.03), RGBColor(0xC0, 0xC8, 0xD4))
    add_text_box(slide1, Inches(1), Inches(4.3), Inches(8), Inches(0.4),
                 f"{len(tasks)} items tracked · SCI-01 through SCI-{max(int(t['code'].split('-')[1]) for t in tasks):02d}",
                 14, False, TEXT_SEC)
    add_aes_logo(slide1)

    # Slide 2 — resume
    slide2 = prs.slides.add_slide(blank)
    add_bg(slide2, BG_LIGHT)
    add_header_bar(slide2)
    set_header_title(slide2, "RESUME")

    kpi_data = [
        (str(len(done)), "Done", GREEN),
        (str(len(ready)), "Ready to move to RPC", PURPLE),
        (str(len(active)), "In Progress / Analysis / Dev", BLUE),
        (str(len(todo)), "To Do", GRAY),
    ]
    card_w = Inches(2.8)
    card_h = Inches(1.8)
    start_x = Inches(0.6)
    gap = Inches(0.25)
    for i, (num, label, color) in enumerate(kpi_data):
        x = start_x + i * (card_w + gap)
        y = Inches(1.4)
        add_rect(slide2, x, y, card_w, card_h, WHITE)
        add_rect(slide2, x, y, card_w, Inches(0.07), color)
        add_text_box(slide2, x, y + Inches(0.25), card_w, Inches(0.9), num, 52, True, color, PP_ALIGN.CENTER)
        add_text_box(slide2, x, y + Inches(1.15), card_w, Inches(0.55), label, 14, True, TEXT_SEC, PP_ALIGN.CENTER)

    badge = add_rounded_rect(slide2, Inches(11.5), Inches(1.5), Inches(1.3), Inches(1), NAVY)
    add_text_box(slide2, Inches(11.5), Inches(1.55), Inches(1.3), Inches(0.6),
                 str(len(tasks)), 32, True, WHITE, PP_ALIGN.CENTER)
    add_text_box(slide2, Inches(11.5), Inches(2.05), Inches(1.3), Inches(0.35),
                 "Total Items", 11, True, CYAN, PP_ALIGN.CENTER)

    ready_codes = ", ".join(t["code"] for t in ready)
    active_codes = ", ".join(t["code"] for t in active)
    todo_codes = ", ".join(t["code"] for t in todo)
    p1_codes = ", ".join(t["code"] for t in priority_items)
    bullets = [
        f"SCI-01 is done ({len(done)} completed item).",
        truncate_words(
            f"Ready to move to RPC ({len(ready)}): {ready_codes}. Deployment target: 02-Jun · CHG1026240.",
            140,
        ),
        truncate_words(
            f"In Progress / Analysis / Development ({len(active)}): {active_codes}.",
            140,
        ),
        truncate_words(f"To Do backlog ({len(todo)}): {todo_codes}.", 140),
        truncate_words(f"Priority 1 focus: {p1_codes}.", 140),
    ]
    add_text_box(slide2, CONTENT_LEFT, Inches(3.45), Inches(4), Inches(0.4), "Overview", 18, True, TEXT_PRI)
    add_multiline_text_box(
        slide2,
        CONTENT_LEFT + Inches(0.2),
        Inches(3.95),
        CONTENT_WIDTH,
        Inches(3.2),
        [f"•  {b}" for b in bullets],
        font_size=12,
        color=TEXT_SEC,
        line_spacing=1.25,
    )

    detail_cols = [
        ("SCI", Inches(0.72)),
        ("Task", Inches(4.85)),
        ("RITM", Inches(1.55)),
        ("Assigned", Inches(1.35)),
        ("Due", Inches(0.72)),
        ("Deploy", Inches(0.72)),
        ("CHG", Inches(1.14)),
    ]

    def ready_rows(task_list):
        return [
            [
                t["code"],
                truncate_words(t["name"], 85),
                truncate_words(t["ritm"], 28),
                truncate_words(t["assigned"], 22),
                t["due"],
                t["deployment"],
                truncate_words(t["change_request"], 18),
            ]
            for t in task_list
        ]

    # Slide 3 — Ready to move to RPC
    add_paginated_table_slides(
        prs,
        blank,
        "READY TO MOVE TO RPC",
        ready_rows(ready),
        detail_cols,
    )

    # Slide 4 — In Progress / Analysis / Development (cards)
    card_row_h = Inches(1.72)
    card_gap_y = Inches(0.12)
    for page_idx, page_tasks in enumerate(chunk(active, 8), start=1):
        slide4 = prs.slides.add_slide(blank)
        add_bg(slide4, BG_LIGHT)
        add_header_bar(slide4)
        suffix = f" ({page_idx})" if len(active) > 8 else ""
        set_header_title(slide4, f"IN PROGRESS — ANALYSIS & DEVELOPMENT{suffix}")
        for i, task in enumerate(page_tasks):
            col = i % 2
            row = i // 2
            x = CONTENT_LEFT + col * Inches(6.35)
            y = TABLE_TOP + row * (card_row_h + card_gap_y)
            accent = BLUE if task["status"] != "In analysis" else CYAN
            detail_card(slide4, x, y, task, accent)

    # Slide 5 — To Do
    todo_cols = [
        ("SCI", Inches(0.72)),
        ("Task", Inches(5.35)),
        ("Assigned", Inches(1.45)),
        ("Requestor", Inches(1.45)),
        ("Due", Inches(0.78)),
    ]
    todo_rows = [
        [
            t["code"],
            truncate_words(t["name"], 95),
            truncate_words(t["assigned"], 24),
            truncate_words(t["requestor"], 24),
            t["due"],
        ]
        for t in todo
    ]
    add_paginated_table_slides(
        prs,
        blank,
        "TO DO",
        todo_rows,
        todo_cols,
        row_height=Inches(0.42),
    )

    # Slide 6 — all tasks by status
    tracker_cols = [
        ("SCI", Inches(0.72)),
        ("Description", Inches(4.35)),
        ("Status", Inches(1.45)),
        ("Type", Inches(1.05)),
        ("P", Inches(0.38)),
        ("Due", Inches(0.72)),
        ("Assigned", Inches(1.25)),
        ("Requestor", Inches(1.23)),
    ]

    ordered_tasks = []
    for status in STATUS_ORDER:
        ordered_tasks.extend([t for t in tasks if t["status"] == status])

    tracker_rows = [
        [
            t["code"],
            truncate_words(t["name"], 78),
            t["status"],
            t["type"],
            t["priority"],
            t["due"],
            truncate_words(t["assigned"], 20),
            truncate_words(t["requestor"], 20),
        ]
        for t in ordered_tasks
    ]

    add_paginated_table_slides(
        prs,
        blank,
        "ALL TASKS — ORDERED BY STATUS",
        tracker_rows,
        tracker_cols,
    )

    # Slide 7 — prioritization
    slide7 = prs.slides.add_slide(blank)
    add_bg(slide7, BG_LIGHT)
    add_header_bar(slide7)
    set_header_title(slide7, "PRIORITIZATION")

    p1_cols = [
        ("SCI", Inches(0.72)),
        ("Task", Inches(3.35)),
        ("Status", Inches(1.35)),
        ("Due", Inches(0.72)),
        ("Assigned", Inches(1.06)),
    ]
    p1_rows = [
        [
            t["code"],
            truncate_words(t["name"], 72),
            t["status"],
            t["due"],
            truncate_words(t["assigned"], 18),
        ]
        for t in priority_items
    ]
    other_cols = [
        ("SCI", Inches(0.72)),
        ("Task", Inches(3.05)),
        ("P", Inches(0.45)),
        ("Status", Inches(1.25)),
        ("Due", Inches(0.72)),
    ]
    other_rows = [
        [
            t["code"],
            truncate_words(t["name"], 68),
            f"P{t['priority']}" if t["priority"] not in ("-", "") else "-",
            t["status"],
            t["due"],
        ]
        for t in tasks
        if t["priority"] != "1"
    ]

    add_text_box(slide7, CONTENT_LEFT, Inches(1.02), Inches(5.5), Inches(0.32), "Priority 1", 15, True, RED)
    p1_table_shape = slide7.shapes.add_table(
        len(p1_rows) + 1,
        len(p1_cols),
        CONTENT_LEFT,
        Inches(1.38),
        sum(w for _, w in p1_cols),
        Inches(0.38) * (len(p1_rows) + 1),
    )
    p1_table = p1_table_shape.table
    for ci, (_, width) in enumerate(p1_cols):
        p1_table.columns[ci].width = width
    for ci, (header, _) in enumerate(p1_cols):
        cell = p1_table.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        set_cell_text(cell, header, font_size=9, bold=True, color=WHITE)
    for ri, row in enumerate(p1_rows, start=1):
        for ci, val in enumerate(row):
            cell = p1_table.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
            color = CYAN if ci == 0 else TEXT_PRI
            bold = ci == 0
            if ci == 2:
                color = STATUS_COLORS.get(str(val), TEXT_PRI)
                bold = True
            elif ci == 3 and str(val) not in ("-", "TBD", ""):
                color = ORANGE
                bold = True
            set_cell_text(cell, val, font_size=8, bold=bold, color=color)

    right_left = Inches(6.85)
    add_text_box(slide7, right_left, Inches(1.02), Inches(5.5), Inches(0.32), "Other priorities", 15, True, BLUE)
    shown_other = other_rows[:14]
    other_table_shape = slide7.shapes.add_table(
        len(shown_other) + 1,
        len(other_cols),
        right_left,
        Inches(1.38),
        sum(w for _, w in other_cols),
        Inches(0.36) * (len(shown_other) + 1),
    )
    other_table = other_table_shape.table
    for ci, (_, width) in enumerate(other_cols):
        other_table.columns[ci].width = width
    for ci, (header, _) in enumerate(other_cols):
        cell = other_table.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = SLATE
        set_cell_text(cell, header, font_size=9, bold=True, color=WHITE)
    for ri, row in enumerate(shown_other, start=1):
        for ci, val in enumerate(row):
            cell = other_table.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
            color = CYAN if ci == 0 else TEXT_PRI
            bold = ci == 0
            if ci == 3:
                color = STATUS_COLORS.get(str(val), TEXT_PRI)
                bold = True
            elif ci == 4 and str(val) not in ("-", "TBD", ""):
                color = ORANGE
                bold = True
            set_cell_text(cell, val, font_size=8, bold=bold, color=color)

    # Slide 8 — thanks
    slide8 = prs.slides.add_slide(blank)
    add_bg(slide8, WHITE)
    add_rect(slide8, Inches(0), Inches(0), Inches(0.08), SLIDE_H, CYAN)
    add_oval(slide8, Inches(9.8), Inches(-1), Inches(4.8), Inches(4.8), RGBColor(0xD6, 0xE4, 0xF7), 30000)
    add_oval(slide8, Inches(10.5), Inches(5.2), Inches(3.5), Inches(3.5), RGBColor(0xD0, 0xF0, 0xF4), 20000)
    add_text_box(slide8, Inches(1), Inches(3), Inches(8), Inches(1), "Thank you", 52, True, NAVY)
    add_rect(slide8, Inches(1), Inches(4.2), Inches(5.5), Inches(0.03), RGBColor(0xC0, 0xC8, 0xD4))
    add_text_box(slide8, Inches(1), Inches(4.5), Inches(8), Inches(0.4),
                 f"Cashflow Project Team · {report_date.strftime('%B %d, %Y')}", 16, False, TEXT_SEC)
    add_aes_logo(slide8, left=Inches(8.15), top=Inches(5.75), width=Inches(3.2))

    prs.save(output_path)


def main():
    base = Path(__file__).resolve().parent
    xlsx_path = base / "CPM-ProcurementForecasting-Status.xlsx"
    report_date = date(2026, 5, 29)
    output_path = base / f"CASHFLOW_Status_{report_date.strftime('%Y%m%d')}.pptx"

    tasks = load_tasks(xlsx_path)
    build_presentation(tasks, report_date, output_path)
    print(f"Saved: {output_path}")
    print(f"Tasks loaded: {len(tasks)}")


if __name__ == "__main__":
    main()
